"""
Gerador de Certificados - Backend Flask
Correcoes:
  - Substituicao de XXXXXX por contexto (seminario, ministrante, mes)
  - Preview via URL publica (sem iframe bloqueado pelo Chrome)
"""

import os, re, io, json, smtplib, zipfile, subprocess, tempfile
from pathlib import Path
from email import encoders
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

import pandas as pd
from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024

UPLOAD_FOLDER = Path(tempfile.gettempdir()) / 'cert_uploads'
OUTPUT_FOLDER = Path(tempfile.gettempdir()) / 'cert_output'
UPLOAD_FOLDER.mkdir(exist_ok=True)
OUTPUT_FOLDER.mkdir(exist_ok=True)


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api/planilha/colunas', methods=['POST'])
def planilha_colunas():
    if 'planilha' not in request.files:
        return jsonify({'erro': 'Nenhum arquivo enviado'}), 400
    f = request.files['planilha']
    ext = f.filename.rsplit('.', 1)[-1].lower()
    try:
        df = pd.read_csv(f, dtype=str).fillna('') if ext == 'csv' \
             else pd.read_excel(f, dtype=str).fillna('')
        df.columns = df.columns.str.strip()
        return jsonify({'colunas': list(df.columns), 'total': len(df),
                        'preview': df.head(5).to_dict(orient='records')})
    except Exception as e:
        return jsonify({'erro': str(e)}), 500


@app.route('/api/preview', methods=['POST'])
def preview():
    """Gera PDF de 1 participante e retorna a URL para abrir em nova aba."""
    if 'template' not in request.files:
        return jsonify({'erro': 'Template nao enviado'}), 400
    template_file = request.files['template']
    try:
        cfg = json.loads(request.form.get('config', '{}'))
    except Exception:
        return jsonify({'erro': 'Config invalida'}), 400

    dados = {
        'nome':        cfg.get('nome', 'Participante Exemplo'),
        'seminario':   cfg.get('seminario', 'Titulo do Seminario'),
        'ministrante': cfg.get('ministrante', 'Nome do Ministrante'),
        'dia':         cfg.get('dia', 'XX'),
        'mes':         cfg.get('mes', 'mes'),
        'ano':         cfg.get('ano', '2026'),
        'carga':       cfg.get('carga', 'uma hora'),
    }

    template_path = UPLOAD_FOLDER / 'preview_template.pptx'
    template_file.save(str(template_path))

    pptx_path = UPLOAD_FOLDER / 'preview.pptx'
    pdf_path  = UPLOAD_FOLDER / 'preview.pdf'
    if pdf_path.exists():
        pdf_path.unlink()

    try:
        gerar_pptx(str(template_path), dados, str(pptx_path))
        converter_pdf(str(pptx_path), str(UPLOAD_FOLDER))
        if not pdf_path.exists():
            return jsonify({'erro': 'LibreOffice nao gerou o PDF'}), 500
        # Retorna URL — o frontend abre em nova aba (sem iframe bloqueado)
        return jsonify({'url': '/api/preview/pdf'})
    except Exception as e:
        return jsonify({'erro': str(e)}), 500


@app.route('/api/preview/pdf')
def preview_pdf():
    """Serve o PDF inline — abre no navegador sem bloqueio."""
    pdf_path = UPLOAD_FOLDER / 'preview.pdf'
    if not pdf_path.exists():
        return 'PDF nao encontrado', 404
    return send_file(str(pdf_path), mimetype='application/pdf',
                     as_attachment=False, download_name='preview.pdf')


@app.route('/api/gerar', methods=['POST'])
def gerar():
    if 'template' not in request.files:
        return jsonify({'erro': 'Template nao enviado'}), 400
    if 'planilha' not in request.files:
        return jsonify({'erro': 'Planilha nao enviada'}), 400

    template_file = request.files['template']
    planilha_file = request.files['planilha']
    try:
        cfg = json.loads(request.form.get('config', '{}'))
    except Exception:
        return jsonify({'erro': 'Config invalida'}), 400

    seminario        = cfg.get('seminario', '')
    ministrante      = cfg.get('ministrante', '')
    dia              = cfg.get('dia', '')
    mes              = cfg.get('mes', '')
    ano              = cfg.get('ano', '2026')
    carga            = cfg.get('carga', 'uma hora')
    col_nome         = cfg.get('col_nome', 'Nome completo')
    col_email        = cfg.get('col_email', '')
    enviar_email_flag = cfg.get('enviar_email', False)
    smtp_host        = cfg.get('smtp_host', 'smtp.gmail.com')
    smtp_port        = int(cfg.get('smtp_port', 587))
    smtp_user        = cfg.get('smtp_user', '')
    smtp_pass        = cfg.get('smtp_pass', '')
    smtp_nome        = cfg.get('smtp_nome', 'Coloquinho da Pos')
    assunto_tpl      = cfg.get('assunto', 'Certificado - {seminario}')
    corpo_tpl        = cfg.get('corpo', CORPO_PADRAO)

    try:
        ext = planilha_file.filename.rsplit('.', 1)[-1].lower()
        df  = pd.read_csv(planilha_file, dtype=str).fillna('') if ext == 'csv' \
              else pd.read_excel(planilha_file, dtype=str).fillna('')
        df.columns = df.columns.str.strip()
    except Exception as e:
        return jsonify({'erro': f'Erro ao ler planilha: {e}'}), 400

    if col_nome not in df.columns:
        return jsonify({'erro': f'Coluna "{col_nome}" nao encontrada. Disponiveis: {list(df.columns)}'}), 400

    template_path = UPLOAD_FOLDER / secure_filename(template_file.filename or 'template.pptx')
    template_file.save(str(template_path))

    for f_ in OUTPUT_FOLDER.glob('*'):
        try: f_.unlink()
        except: pass

    resultados = []

    for idx, row in df.iterrows():
        nome = str(row.get(col_nome, '')).strip()
        if not nome:
            resultados.append({'idx': idx+1, 'nome': '(vazio)', 'status': 'pulado', 'msg': 'Nome vazio'})
            continue

        email_dest = str(row.get(col_email, '')).strip() if col_email else ''
        dados = {'nome': nome, 'seminario': seminario, 'ministrante': ministrante,
                 'dia': dia, 'mes': mes, 'ano': ano, 'carga': carga}

        try:
            nome_arq  = nome_seguro(nome)
            pptx_path = OUTPUT_FOLDER / f'{nome_arq}.pptx'
            pdf_path  = OUTPUT_FOLDER / f'{nome_arq}.pdf'

            gerar_pptx(str(template_path), dados, str(pptx_path))
            converter_pdf(str(pptx_path), str(OUTPUT_FOLDER))

            if not pdf_path.exists():
                raise FileNotFoundError('PDF nao gerado')

            msg_email = ''
            if enviar_email_flag and email_dest and smtp_user and smtp_pass:
                subst = {'{nome}': nome, '{seminario}': seminario,
                         '{data}': f'{dia} de {mes} de {ano}', '{ministrante}': ministrante}
                enviar_email(smtp_host, smtp_port, smtp_user, smtp_pass, smtp_nome,
                             email_dest,
                             substituir_dict(assunto_tpl, subst),
                             substituir_dict(corpo_tpl,   subst),
                             str(pdf_path))
                msg_email = f' -> {email_dest}'

            resultados.append({'idx': idx+1, 'nome': nome, 'status': 'ok',
                                'arquivo': nome_arq + '.pdf', 'msg': 'Gerado' + msg_email})
        except Exception as e:
            resultados.append({'idx': idx+1, 'nome': nome, 'status': 'erro', 'msg': str(e)})

    return jsonify({
        'resultados': resultados, 'total': len(df),
        'sucesso': sum(1 for r in resultados if r['status'] == 'ok'),
        'erros':   sum(1 for r in resultados if r['status'] == 'erro'),
    })


@app.route('/api/download-zip')
def download_zip():
    pdfs = list(OUTPUT_FOLDER.glob('*.pdf'))
    if not pdfs:
        return jsonify({'erro': 'Nenhum PDF encontrado'}), 404
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for pdf in pdfs:
            zf.write(pdf, pdf.name)
    buf.seek(0)
    return send_file(buf, mimetype='application/zip',
                     as_attachment=True, download_name='certificados.zip')


# ══ FUNCOES ══

def gerar_pptx(template_path: str, dados: dict, saida: str):
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                texto = ''.join(r.text for r in para.runs)
                if not texto.strip():
                    continue
                texto_novo = substituir_marcadores(texto, dados)
                if texto != texto_novo and para.runs:
                    para.runs[0].text = texto_novo
                    for r in para.runs[1:]:
                        r.text = ''
    prs.save(saida)


def substituir_marcadores(texto: str, d: dict) -> str:
    """
    Substitui marcadores do template Coloquinho da Pos por contexto.

    Texto reconstruido dos runs:
    'Certificamos que <<Nome completo>> frequentou o seminario "XXXXXX",
     no Coloquinho da Pos no IME-USP, ministrado por XXXXXX
     no dia XX de XXXXXX de 2026, com carga horaria de uma hora.'

    Os tres XXXXXX sao distintos pelo contexto ao redor.
    """
    resultado = texto

    # 1. Nome do participante
    resultado = resultado.replace('<<Nome completo>>', d.get('nome', ''))

    # 2. Data: "XX de XXXXXX de 2026" -> "15 de marco de 2026"
    resultado = re.sub(
        r'\bXX\b(\s+de\s+)XXXXXX(\s+de\s+\d{4})',
        d.get('dia', 'XX') + r'\g<1>' + d.get('mes', 'XXXXXX') + r'\g<2>',
        resultado
    )

    # 3. Seminario: contexto 'seminario "XXXXXX"'
    #    Aspas tipograficas (\u201c \u201d) ou retas
    resultado = re.sub(
        r'(semin[a\u00e1]rio\s+[\u201c\"])XXXXXX([\u201d\"])',
        lambda m: m.group(1) + d.get('seminario', 'XXXXXX') + m.group(2),
        resultado
    )

    # 4. Ministrante: contexto 'ministrado por XXXXXX'
    resultado = re.sub(
        r'(ministrado por\s+)XXXXXX',
        lambda m: m.group(1) + d.get('ministrante', 'XXXXXX'),
        resultado
    )

    # 5. Carga horaria (opcional, se diferente do padrao)
    if d.get('carga') and d['carga'] != 'uma hora':
        resultado = re.sub(
            r'(carga hor[a\u00e1]ria de\s+)uma hora',
            lambda m: m.group(1) + d['carga'],
            resultado
        )

    return resultado


def converter_pdf(pptx_path: str, pasta_saida: str):
    res = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pdf',
         '--outdir', pasta_saida, pptx_path],
        capture_output=True, text=True, timeout=120
    )
    if res.returncode != 0:
        raise RuntimeError(f'LibreOffice: {res.stderr[:400]}')


def enviar_email(host, port, usuario, senha, nome_rem, dest, assunto, corpo_html, pdf_path):
    msg = MIMEMultipart('alternative')
    msg['Subject'] = assunto
    msg['From']    = f'{nome_rem} <{usuario}>'
    msg['To']      = dest
    msg.attach(MIMEText(re.sub(r'<[^>]+>', '', corpo_html), 'plain', 'utf-8'))
    msg.attach(MIMEText(corpo_html, 'html', 'utf-8'))
    with open(pdf_path, 'rb') as f:
        parte = MIMEBase('application', 'pdf')
        parte.set_payload(f.read())
    encoders.encode_base64(parte)
    parte.add_header('Content-Disposition', 'attachment', filename=Path(pdf_path).name)
    msg.attach(parte)
    with smtplib.SMTP(host, port) as smtp:
        smtp.starttls()
        smtp.login(usuario, senha)
        smtp.sendmail(usuario, dest, msg.as_string())


def substituir_dict(texto, d):
    for k, v in d.items():
        texto = texto.replace(k, v)
    return texto


def nome_seguro(nome: str) -> str:
    return re.sub(r'[<>:"/\\|?*\x00-\x1f]', '', nome).strip() or 'certificado'


CORPO_PADRAO = (
    "<html><body style='font-family:Arial,sans-serif;max-width:600px;margin:0 auto;'>"
    "<div style='background:#1a1410;color:#f7f3ec;padding:28px;border-radius:8px 8px 0 0;'>"
    "<h2 style='margin:0;'>Seu certificado chegou!</h2></div>"
    "<div style='background:#f9f6f0;padding:28px;border:1px solid #e0d8cc;border-radius:0 0 8px 8px;'>"
    "<p>Ola, <strong>{nome}</strong>!</p>"
    "<p>Segue em anexo o seu certificado de participacao no seminario "
    "<strong>\"{seminario}\"</strong>, realizado no Coloquinho da Pos no IME-USP em {data}.</p>"
    "<p>Obrigado pela sua presenca!</p>"
    "<hr style='border:none;border-top:1px solid #e0d8cc;margin:20px 0;'>"
    "<p style='color:#999;font-size:12px;'>Organizacao: Raquel Mansano Goncalves Cenciarelli</p>"
    "</div></body></html>"
)

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)

